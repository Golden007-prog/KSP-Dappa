"""Rasterise the built deck to PNGs so the layout can be eyeballed without PowerPoint.

Not a pixel-perfect renderer — it draws the background picture, then the text
boxes and images at their real positions with their real colours. Enough to catch
overflow, invisible text and misplaced images.

    python3.12 scripts/preview_deck.py
"""
import io
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
DATATHON = os.path.dirname(REPO)
DECK = os.path.join(DATATHON, "KSP Datathon 2026 _ Prototype Submission - KSP DAPPA.pptx")
OUT = os.path.join(REPO, "docs", "screenshots", "deck-preview")

SCALE = 130  # px per inch
EMU_IN = 914400


def font(size_pt, bold=False):
    names = ["seguisb.ttf" if bold else "segoeui.ttf", "arialbd.ttf" if bold else "arial.ttf"]
    for n in names:
        for d in (r"C:\Windows\Fonts", "/usr/share/fonts/truetype/dejavu"):
            p = os.path.join(d, n)
            if os.path.exists(p):
                try:
                    return ImageFont.truetype(p, int(size_pt * SCALE / 72))
                except OSError:
                    pass
    return ImageFont.load_default()


def wrap(draw, text, fnt, max_px):
    words, lines, cur = text.split(), [], ""
    for w in words:
        trial = (cur + " " + w).strip()
        if draw.textlength(trial, font=fnt) <= max_px or not cur:
            cur = trial
        else:
            lines.append(cur)
            cur = w
    if cur:
        lines.append(cur)
    return lines


def main():
    os.makedirs(OUT, exist_ok=True)
    prs = Presentation(DECK)
    W = int(prs.slide_width / EMU_IN * SCALE)
    H = int(prs.slide_height / EMU_IN * SCALE)

    for i, slide in enumerate(prs.slides, start=1):
        canvas = Image.new("RGB", (W, H), (0, 0, 0))
        # backgrounds first, then foreground images
        pics = [s for s in slide.shapes if s.shape_type == 13]
        for sh in sorted(pics, key=lambda s: -(s.width or 0) * (s.height or 0)):
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
            except Exception:
                continue
            w = int((sh.width or 0) / EMU_IN * SCALE)
            h = int((sh.height or 0) / EMU_IN * SCALE)
            if w <= 0 or h <= 0:
                continue
            canvas.paste(im.resize((w, h), Image.LANCZOS),
                         (int((sh.left or 0) / EMU_IN * SCALE),
                          int((sh.top or 0) / EMU_IN * SCALE)))

        d = ImageDraw.Draw(canvas)
        for sh in slide.shapes:
            if sh.shape_type == 13 or not sh.has_text_frame:
                continue
            x = int((sh.left or 0) / EMU_IN * SCALE)
            y = int((sh.top or 0) / EMU_IN * SCALE)
            maxw = int((sh.width or 0) / EMU_IN * SCALE) - 6
            for para in sh.text_frame.paragraphs:
                runs = para.runs
                if not runs:
                    y += 8
                    continue
                text = "".join(r.text for r in runs)
                r0 = runs[0]
                size = (r0.font.size.pt if r0.font.size else 13.5)
                bold = bool(r0.font.bold)
                try:
                    col = r0.font.color.rgb
                    fill = (col[0], col[1], col[2]) if col else (236, 241, 248)
                except Exception:
                    fill = (236, 241, 248)
                fnt = font(size, bold)
                indent = 22 if para.level else 0
                for line in wrap(d, text, fnt, maxw - indent):
                    d.text((x + indent, y), line, font=fnt, fill=fill)
                    y += int(size * SCALE / 72 * 1.28)
                y += 5

        path = os.path.join(OUT, f"slide-{i:02d}.png")
        canvas.save(path)
    print("wrote", len(prs.slides.__iter__.__self__._sldIdLst), "previews to", OUT)


if __name__ == "__main__":
    main()
