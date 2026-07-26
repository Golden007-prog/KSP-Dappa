"""Fill the KSP Datathon 2026 Prototype Submission Template with KSP DAPPA content.

The template ships 16 slides, each a full-bleed background picture plus a heading
text box. This script leaves the template chrome untouched and adds a body block
(and, where useful, an image) under each heading.

Content lives in deck_content.py so the copy can be reviewed on its own.

    python3.12 scripts/build_submission_deck.py

Writes  "KSP Datathon 2026 _ Prototype Submission - KSP DAPPA.pptx"  next to the
template, leaving the original template file untouched.
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(HERE)
DATATHON = os.path.dirname(REPO)
SHOTS = os.path.join(REPO, "docs", "screenshots")

TEMPLATE = os.path.join(DATATHON, "KSP Datathon 2026 _ Prototype Submission Template.pptx")
OUTPUT = os.path.join(DATATHON, "KSP Datathon 2026 _ Prototype Submission - KSP DAPPA.pptx")

sys.path.insert(0, HERE)
from deck_content import SLIDES, TEAM  # noqa: E402

# The template ships a full-bleed BLACK background on every content slide, so
# body copy has to be light. Sampled from the template art, not assumed.
INK = RGBColor(0xEC, 0xF1, 0xF8)
MUTED = RGBColor(0xA8, 0xB6, 0xCC)
ACCENT = RGBColor(0xF0, 0xB4, 0x29)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(10.0)


def heading_bottom(slide, default=1.55):
    """Lowest edge of the template's heading text box, in inches."""
    best = None
    for sh in slide.shapes:
        if sh.shape_type == 13:  # PICTURE -> the full-bleed background
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            bottom = (sh.top + sh.height) / 914400
            best = bottom if best is None else max(best, bottom)
    return best if best else default


def add_bullets(slide, top_in, bullets, sub, left_in=0.55, width_in=8.9,
                size=13.5, sub_size=11.5):
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in),
        Inches(5.62 - top_in - 0.25))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    children = {}
    for entry in sub or []:
        idx, _, text = entry.partition("|")
        try:
            children.setdefault(int(idx), []).append(text.strip())
        except ValueError:
            continue

    first = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        marker = p.add_run()
        marker.text = "▪  "
        marker.font.size = Pt(size)
        marker.font.color.rgb = ACCENT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        run.font.bold = False
        for child in children.get(i, []):
            cp = tf.add_paragraph()
            cp.space_after = Pt(3)
            cp.level = 1
            cr = cp.add_run()
            cr.text = "–  " + child
            cr.font.size = Pt(sub_size)
            cr.font.color.rgb = MUTED
    return box


def add_image(slide, path, left_in, top_in, max_w_in, max_h_in, border=True):
    """Place an image fitted inside a box, preserving aspect ratio."""
    if not os.path.exists(path):
        print(f"  ! missing image {path}")
        return None
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w_in / (iw / 96), max_h_in / (ih / 96))
    w = (iw / 96) * scale
    h = (ih / 96) * scale
    left = Inches(left_in + (max_w_in - w) / 2)
    top = Inches(top_in + (max_h_in - h) / 2)
    pic = slide.shapes.add_picture(path, left, top, Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = RGBColor(0x25, 0x33, 0x4D)
        pic.line.width = Pt(0.75)
    return pic


def caption(slide, text, top_in, left_in=0.55, width_in=8.9, size=9.5):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(0.28))
    p = box.text_frame.paragraphs[0]
    box.text_frame.word_wrap = True
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.CENTER
    return box


def fill_team(slide):
    """Slide 1 already carries the 'Team Details' labels; replace them wholesale."""
    for sh in slide.shapes:
        if sh.has_text_frame and "Team Details" in sh.text_frame.text:
            tf = sh.text_frame
            tf.clear()
            rows = [
                ("Team Details", None),
                ("Team name", TEAM["team_name"]),
                ("Team leader name", TEAM["leader"]),
                ("Team size", TEAM["size"]),
                ("Problem Statement", TEAM["problem"]),
            ]
            first = True
            for label, value in rows:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.space_after = Pt(6)
                r = p.add_run()
                if value is None:
                    r.text = label
                    r.font.size = Pt(20)
                    r.font.bold = True
                else:
                    r.text = f"{label}:  "
                    r.font.size = Pt(13)
                    r.font.bold = True
                    v = p.add_run()
                    v.text = value
                    v.font.size = Pt(13)
                    v.font.bold = False
            return sh
    return None


def main():
    prs = Presentation(TEMPLATE)
    by_number = {s["slideNumber"]: s for s in SLIDES}

    for idx, slide in enumerate(prs.slides, start=1):
        spec = by_number.get(idx)
        if idx == 1:
            fill_team(slide)
            print("slide 1  team details")
            continue
        if spec is None:
            continue

        top = heading_bottom(slide) + 0.12
        images = spec.get("images") or []
        bullets = spec.get("bullets") or []

        if images and bullets:
            # bullets on the left, image on the right
            add_bullets(slide, top, bullets, spec.get("subBullets"),
                        left_in=0.45, width_in=5.05,
                        size=spec.get("size", 11.5), sub_size=10.0)
            add_image(slide, os.path.join(SHOTS, images[0]),
                      5.62, top, 3.95, 5.62 - top - 0.32)
            if spec.get("caption"):
                caption(slide, spec["caption"], 5.62 - 0.30, left_in=5.62, width_in=3.95, size=8.5)
        elif images:
            n = len(images)
            if n == 1:
                add_image(slide, os.path.join(SHOTS, images[0]),
                          0.45, top, 9.1, 5.62 - top - 0.34)
            else:
                cols = 2 if n <= 4 else 3
                rows = (n + cols - 1) // cols
                cw = 9.1 / cols
                ch = (5.62 - top - 0.34) / rows
                for i, img in enumerate(images):
                    r, c = divmod(i, cols)
                    add_image(slide, os.path.join(SHOTS, img),
                              0.45 + c * cw, top + r * ch, cw - 0.12, ch - 0.10)
            if spec.get("caption"):
                caption(slide, spec["caption"], 5.62 - 0.30)
        else:
            add_bullets(slide, top, bullets, spec.get("subBullets"),
                        size=spec.get("size", 13.5))

        print(f"slide {idx:<2} {len(bullets)} bullets, {len(images)} image(s)")

    prs.save(OUTPUT)
    print("\nwrote", OUTPUT)


if __name__ == "__main__":
    main()
